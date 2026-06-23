"""
title: PowerPoint Creator
author: open-webui-aws-fargate
version: 2.0.0
description: Generate a Microsoft PowerPoint (.pptx) presentation from a Markdown string. Use # headings to delimit slides; bullet lines become slide bullets; other text becomes body paragraphs. The file is delivered as a chat attachment.
required_open_webui_version: 0.9.0
"""

import io
import re
import uuid
from typing import Optional

from pydantic import BaseModel, Field

from pptx import Presentation
from pptx.util import Pt


class Tools:
    class Valves(BaseModel):
        max_slides: int = Field(default=100, description="Refuse decks with more than this many slides.")
        max_size_mb: int = Field(default=20, description="Refuse to attach files larger than this many megabytes.")

    def __init__(self):
        self.valves = self.Valves()
        self.citation = False

    async def create_powerpoint(
        self,
        title: str,
        markdown: str,
        __event_emitter__=None,
        __user__: Optional[dict] = None,
    ) -> str:
        """
        Create a Microsoft PowerPoint (.pptx) deck and attach it to the chat.

        Args:
            title: Deck title. Used for the title slide and the filename.
            markdown: Slide content as Markdown. Formatting rules:
                - # Slide Title  starts a new slide with that title
                - Lines starting with - or * become bullet points on the current slide
                - Lines starting with 1. or 1) become numbered bullets
                - Other non-blank lines become body paragraph text
                - ## or ### headings within a slide are rendered as bold body text
                Write one # heading per slide. All content below it belongs to that slide
                until the next # heading.
        """
        slides = _parse_slides(markdown)

        if len(slides) > self.valves.max_slides:
            return f"Error: too many slides ({len(slides)}); limit is {self.valves.max_slides}."

        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        # Title slide
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = ""

        for entry in slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = entry["title"]

            body_ph = next(
                (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
                None,
            )
            if body_ph is None:
                continue

            tf = body_ph.text_frame
            tf.clear()
            first = True

            for item in entry["items"]:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                para.text = item["text"]
                if item["kind"] == "bullet":
                    para.level = 0
                elif item["kind"] == "numbered":
                    para.level = 0
                else:
                    # Body text: indent one level and slightly smaller
                    para.level = 1
                    if para.runs:
                        para.runs[0].font.size = Pt(14)

        buffer = io.BytesIO()
        prs.save(buffer)
        data = buffer.getvalue()

        return await _attach_file(
            __event_emitter__,
            __user__,
            data=data,
            filename=f"{_safe_filename(title)}.pptx",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            max_size_mb=self.valves.max_size_mb,
        )


def _parse_slides(markdown: str) -> list[dict]:
    """
    Split markdown into slides.

    Each # heading starts a new slide. Returns:
      [{"title": str, "items": [{"kind": bullet|numbered|body, "text": str}]}]
    """
    slides = []
    current: Optional[dict] = None

    for line in markdown.splitlines():
        stripped = line.strip()

        # New slide on # heading (exactly one #)
        m = re.match(r"^#(?!#)\s+(.*)", stripped)
        if m:
            if current is not None:
                slides.append(current)
            current = {"title": m.group(1).strip(), "items": []}
            continue

        if current is None or not stripped:
            continue

        # Bullet
        m = re.match(r"^[-*•]\s+(.*)", stripped)
        if m:
            current["items"].append({"kind": "bullet", "text": m.group(1).strip()})
            continue

        # Numbered
        m = re.match(r"^\d+[.)]\s+(.*)", stripped)
        if m:
            current["items"].append({"kind": "numbered", "text": m.group(0).strip()})
            continue

        # ## or ### sub-heading → bold body text
        m = re.match(r"^#{2,6}\s+(.*)", stripped)
        if m:
            current["items"].append({"kind": "body", "text": m.group(1).strip()})
            continue

        # Plain body line
        current["items"].append({"kind": "body", "text": stripped})

    if current is not None:
        slides.append(current)

    return slides


def _safe_filename(name: str, fallback: str = "presentation") -> str:
    name = re.sub(r"[^\w\-. ]", "_", name).strip()
    return name[:120] or fallback


async def _attach_file(
    event_emitter,
    user: Optional[dict],
    *,
    data: bytes,
    filename: str,
    content_type: str,
    max_size_mb: int = 10,
) -> str:
    if event_emitter is None:
        return f"(no event emitter; would attach {filename}, {len(data)} bytes)"
    if len(data) > max_size_mb * 1024 * 1024:
        return f"Refusing to attach: {filename} is {len(data) / 1e6:.1f} MB (limit {max_size_mb} MB)."

    from open_webui.models.files import Files, FileForm
    from open_webui.storage.provider import Storage

    file_id = str(uuid.uuid4())
    user_id = (user or {}).get("id", "")
    storage_filename = f"{file_id}_{filename}"

    _, storage_path = Storage.upload_file(io.BytesIO(data), storage_filename, {})

    record = await Files.insert_new_file(
        user_id,
        FileForm(
            id=file_id,
            filename=filename,
            path=storage_path,
            meta={"name": filename, "content_type": content_type, "size": len(data)},
        ),
    )
    if record is None:
        return f"Failed to register file {filename} with Open WebUI."

    await event_emitter(
        {
            "type": "files",
            "data": {
                "files": [
                    {
                        "type": "file",
                        "id": record.id,
                        "url": f"/api/v1/files/{record.id}",
                        "name": filename,
                        "size": len(data),
                        "status": "uploaded",
                        "error": "",
                        "itemId": str(uuid.uuid4()),
                    }
                ],
            },
        }
    )

    return f"Created **{filename}** ({len(data):,} bytes). The file is now attached to the chat."
